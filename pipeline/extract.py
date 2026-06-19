# -*- coding: utf-8 -*-
"""Extract text from PPTX/DOCX/MD with location markers into per-file dumps."""
import os, sys, glob, io

BASE = r"c:\Users\Chrisi\Downloads\vizerektor-demo\material"
OUT  = r"c:\tmp\vizerektor-extraktion"

def w(fh, s):
    fh.write(s + "\n")

def extract_pptx(path, out):
    from pptx import Presentation
    prs = Presentation(path)
    with io.open(out, "w", encoding="utf-8") as fh:
        for i, slide in enumerate(prs.slides, start=1):
            w(fh, "\n===== FOLIE %d =====" % i)
            # title if present
            for shape in slide.shapes:
                if shape.has_text_frame:
                    txt = shape.text_frame.text.strip()
                    if txt:
                        for line in txt.splitlines():
                            line = line.strip()
                            if line:
                                w(fh, line)
                if shape.has_table:
                    tbl = shape.table
                    for r in tbl.rows:
                        cells = [c.text.strip() for c in r.cells]
                        w(fh, " | ".join(cells))
            # notes
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    w(fh, "[NOTIZEN] " + notes.replace("\n", " "))

def extract_docx(path, out):
    import docx
    d = docx.Document(path)
    with io.open(out, "w", encoding="utf-8") as fh:
        for i, p in enumerate(d.paragraphs, start=1):
            t = p.text.strip()
            if t:
                w(fh, "[P%d] %s" % (i, t))
        for ti, tbl in enumerate(d.tables, start=1):
            w(fh, "\n[TABELLE %d]" % ti)
            for r in tbl.rows:
                cells = [c.text.strip() for c in r.cells]
                w(fh, " | ".join(cells))

def copy_md(path, out):
    with io.open(path, "r", encoding="utf-8") as f:
        lines = f.readlines()
    with io.open(out, "w", encoding="utf-8") as fh:
        for i, line in enumerate(lines, start=1):
            fh.write("[Z%d] %s" % (i, line if line.endswith("\n") else line+"\n"))

files = []
for pat in ("folien/*.pptx", "folien/*.docx", "wissensdokumente/*.md"):
    files += glob.glob(os.path.join(BASE, pat))

manifest = []
for path in sorted(files):
    name = os.path.basename(path)
    ext = name.lower().rsplit(".",1)[-1]
    out = os.path.join(OUT, name + ".txt")
    try:
        if ext == "pptx":
            extract_pptx(path, out)
        elif ext == "docx":
            extract_docx(path, out)
        elif ext == "md":
            copy_md(path, out)
        size = os.path.getsize(out)
        with io.open(out, encoding="utf-8") as f:
            nlines = sum(1 for _ in f)
        manifest.append((name, size, nlines))
        print("OK  %-55s %8d bytes  %5d lines" % (name, size, nlines))
    except Exception as e:
        print("ERR %-55s %s" % (name, e))

total = sum(m[1] for m in manifest)
print("\nTOTAL extracted text: %d bytes (~%d k tokens rough)" % (total, total//4//1000))
